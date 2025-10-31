"""
Version: 1.4
Author: ChatGPT
Description: 企业级PPT生成，带IBM Logo和企业风格背景
Date: 2025-10-31
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# =========================
# 创建PPT对象
# =========================
pres = Presentation()

# =========================
# 全局样式参数
# =========================
TITLE_COLOR = RGBColor(255, 255, 255)      # 白色标题
SUBTITLE_COLOR = RGBColor(200, 200, 200)   # 灰色副标题
BODY_COLOR = RGBColor(255, 255, 255)       # 白色正文
FONT_NAME = 'Arial'
BACKGROUND_COLOR = RGBColor(0, 51, 102)    # IBM深蓝色企业背景
LOGO_PATH = 'ibm_logo.png'                 # IBM logo 图片路径，本地图片

# =========================
# 背景填充函数
# =========================
def set_background(slide, color=BACKGROUND_COLOR):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# =========================
# 添加企业Logo
# =========================

LOGO_PATH = r'C:\\Users\\86181\\Desktop\\Vscode_python\\resources\\ibm.png'

def add_logo(slide, logo_path=LOGO_PATH):
    """在右上角添加Logo"""
    slide.shapes.add_picture(logo_path, Inches(8), Inches(0.2), width=Inches(1.2))

# =========================
# 第一页：标题页
# =========================
slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(slide_layout)
set_background(slide)
add_logo(slide)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Welcome to Python"
subtitle.text = "Life is short, I use Python"

title_tf = title.text_frame.paragraphs[0]
title_tf.font.size = Pt(48)
title_tf.font.bold = True
title_tf.font.color.rgb = TITLE_COLOR
title_tf.font.name = FONT_NAME

subtitle_tf = subtitle.text_frame.paragraphs[0]
subtitle_tf.font.size = Pt(28)
subtitle_tf.font.italic = True
subtitle_tf.font.color.rgb = SUBTITLE_COLOR
subtitle_tf.font.name = FONT_NAME

# =========================
# 第二页：内容页
# =========================
slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(slide_layout)
set_background(slide)
add_logo(slide)

title_shape = slide.shapes.title
title_shape.text = "Introduction"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
title_shape.text_frame.paragraphs[0].font.name = FONT_NAME

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.clear()

contents = [
    ("History of Python", 0),
    ("X'mas 1989", 1),
    ("Guido began to write interpreter for Python.", 2),
]

for text, level in contents:
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.name = FONT_NAME
    if level == 0:
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BODY_COLOR
    elif level == 1:
        p.font.size = Pt(24)
        p.font.color.rgb = BODY_COLOR
    else:
        p.font.size = Pt(22)
        p.font.color.rgb = BODY_COLOR
    p.space_after = Pt(5)

# =========================
# 第三页：图标内容页
# =========================
slide_layout = pres.slide_layouts[5]  # 空白页
slide = pres.slides.add_slide(slide_layout)
set_background(slide)
add_logo(slide)

# 标题
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_shape_tf = title_shape.text_frame
title_shape_tf.text = "Python Features"
title_shape_tf.paragraphs[0].font.size = Pt(36)
title_shape_tf.paragraphs[0].font.bold = True
title_shape_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

# 图标+文字
features = [
    "Easy to Learn",
    "Open Source",
    "Versatile",
    "Large Community"
]
top = 1.5
for feature in features:
    # 圆形图标
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(top), Inches(0.3), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.color.rgb = TITLE_COLOR

    # 文字
    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(8), Inches(0.3))
    p = textbox.text_frame.paragraphs[0]
    p.text = feature
    p.font.size = Pt(24)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_NAME

    top += 0.6

# =========================
# 保存PPT
# =========================
pres.save('test_enterprise_ibm.pptx')
print("企业级PPT已生成：test_enterprise_ibm.pptx")
