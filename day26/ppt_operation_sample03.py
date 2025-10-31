"""
Version: 1.2
Author: ChatGPT
Description: 高级美化PPT生成
Date: 2025-10-31
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# =========================
# 创建PPT对象
# =========================
pres = Presentation()

# =========================
# 全局样式参数
# =========================
TITLE_COLOR = RGBColor(0, 102, 204)      # 深蓝
SUBTITLE_COLOR = RGBColor(102, 102, 102) # 灰色
BODY_COLOR = RGBColor(0, 51, 102)        # 深色正文
FONT_NAME = 'Arial'

# =========================
# 第一页：标题页
# =========================
slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Welcome to Python"
subtitle.text = "Life is short, I use Python"

# 美化标题
title_tf = title.text_frame.paragraphs[0]
title_tf.font.size = Pt(48)
title_tf.font.bold = True
title_tf.font.color.rgb = TITLE_COLOR
title_tf.font.name = FONT_NAME

subtitle_tf = subtitle.text_frame.paragraphs[0]
subtitle_tf.font.size = Pt(28)
subtitle_tf.font.italic = True
subtitle_tf.font.color.rgb = SUBTITLE_COLOR
subtitle_tf.font.name = FONT_NAME

# =========================
# 第二页：内容页
# =========================
slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(slide_layout)

# 设置标题
title_shape = slide.shapes.title
title_shape.text = "Introduction"
title_shape.text_frame.paragraphs[0].font.size = Pt(36)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
title_shape.text_frame.paragraphs[0].font.name = FONT_NAME

# 设置主体内容
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.clear()  # 清空默认文本

# 一级段落
p1 = tf.add_paragraph()
p1.text = 'History of Python'
p1.level = 0
p1.font.size = Pt(28)
p1.font.bold = True
p1.font.color.rgb = BODY_COLOR
p1.font.name = FONT_NAME
p1.space_after = Pt(10)

# 二级段落
p2 = tf.add_paragraph()
p2.text = "X'mas 1989"
p2.level = 1
p2.font.size = Pt(24)
p2.font.color.rgb = BODY_COLOR
p2.font.name = FONT_NAME
p2.space_after = Pt(5)

# 三级段落
p3 = tf.add_paragraph()
p3.text = "Guido began to write interpreter for Python."
p3.level = 2
p3.font.size = Pt(22)
p3.font.color.rgb = BODY_COLOR
p3.font.name = FONT_NAME
p3.space_after = Pt(5)

# =========================
# 第三页：带图标的内容页示例
# =========================
slide_layout = pres.slide_layouts[5]  # 空白页
slide = pres.slides.add_slide(slide_layout)

# 添加标题
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_shape_tf = title_shape.text_frame
title_shape_tf.text = "Python Features"
title_shape_tf.paragraphs[0].font.size = Pt(36)
title_shape_tf.paragraphs[0].font.bold = True
title_shape_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

# 添加图标+文本
features = [
    "Easy to Learn",
    "Open Source",
    "Versatile",
    "Large Community"
]

top = 1.5
for feature in features:
    # 添加圆形图标
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(top), Inches(0.3), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.color.rgb = TITLE_COLOR

    # 添加文字
    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(8), Inches(0.3))
    p = textbox.text_frame.paragraphs[0]
    p.text = feature
    p.font.size = Pt(24)
    p.font.color.rgb = BODY_COLOR
    p.font.name = FONT_NAME

    top += 0.6

# =========================
# 保存PPT
# =========================
pres.save('test_beautiful.pptx')
print("PPT 已生成：test_beautiful.pptx")
