"""
Version: 1.1
Author: ChatGPT
Description: 优化生成PowerPoint
Date: 2025-10-31
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# 创建幻灯片对象
pres = Presentation()

# =========================
# 第一页：标题页
# =========================
title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)

# 设置标题和副标题
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Welcome to Python"
subtitle.text = "Life is short, I use Python"

# 美化标题
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)  # 蓝色

subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.italic = True
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)  # 灰色

# =========================
# 第二页：内容页
# =========================
bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)

# 设置标题
title_shape = slide.shapes.title
title_shape.text = "Introduction"
title_shape.text_frame.paragraphs[0].font.size = Pt(32)
title_shape.text_frame.paragraphs[0].font.bold = True
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# 设置主体内容
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.clear()  # 清空默认文本

# 一级段落
p = tf.add_paragraph()
p.text = 'History of Python'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)
p.level = 0

# 二级段落
p = tf.add_paragraph()
p.text = "X'mas 1989"
p.font.size = Pt(24)
p.level = 1
p.space_after = Pt(5)  # 段落间距

p = tf.add_paragraph()
p.text = "Guido began to write interpreter for Python."
p.font.size = Pt(24)
p.level = 2
p.space_after = Pt(5)

# =========================
# 保存幻灯片
# =========================
pres.save('test_optimized.pptx')
print("PPT 已生成：test_optimized.pptx")
